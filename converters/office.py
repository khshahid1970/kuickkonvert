"""Office <-> PDF conversions, powered by headless LibreOffice (plus a
pdf2docx assist for PDF -> Word, and an image-slide fallback for PDF -> PPT
-- see convert_pdf_to_word and convert_pdf_to_ppt below).

LibreOffice is the only free/open-source engine that reliably round-trips
Word/Excel/PowerPoint <-> PDF with acceptable layout fidelity. Each call
shells out to `soffice --headless --convert-to`, which is CPU-bound and
single-document-at-a-time by nature of the LO profile lock, so callers
should expect a real conversion (not instant) and should not call this
concurrently against the same user profile directory.
"""
import glob
import io
import os
import subprocess

SOFFICE_BIN = os.environ.get("SOFFICE_BIN", "soffice")
CONVERT_TIMEOUT = int(os.environ.get("CONVERT_TIMEOUT_SECONDS", "120"))

SPREADSHEET_EXTS = {"xls", "xlsx"}

# LibreOffice's Calc PDF-export filter accepts extra filter data as JSON
# appended to the filter name. SinglePageSheets forces every worksheet onto
# one PDF page, growing the page size to fit the data instead of Calc's
# default behaviour of paginating a wide/tall sheet across several separate
# pages -- which is what was cutting off columns and rows entirely in
# "Excel to PDF" conversions. Combined with _widen_columns_to_fit() below
# (which fixes text getting visually clipped within a column that's too
# narrow for its content), this makes the output PDF's page dimensions
# track the actual size of the input data rather than an arbitrary fixed
# page size. Verified against LibreOffice 24.2 in testing on 2026-08-27.
_SINGLE_PAGE_SHEETS_FILTER = (
    'pdf:calc_pdf_Export:{"SinglePageSheets":{"type":"boolean","value":"true"}}'
)


class ConversionError(Exception):
    pass


def _run_soffice(input_path: str, out_dir: str, target_filter: str, infilter: str = None):
    """Run soffice --convert-to and return the produced file path.

    Output is written to a dedicated subdirectory, never the same directory
    the input file lives in -- if input and output shared a directory, an
    input like "report.pdf" being converted to "report.pdf"-named output
    (or a glob matching both the original upload and the new file when they
    share a stem) could get confused for one another and the wrong file
    would be returned to the user.

    `infilter` forces how LibreOffice interprets the *input* file. This
    matters specifically for PDF input: LibreOffice opens a PDF as a Draw
    (graphics) document by default, which can export fine to PPTX (another
    graphics-ish format) but fails outright when asked to save as DOCX --
    Draw has no Writer text model to save from. Passing
    infilter="writer_pdf_import" makes LibreOffice reimport the PDF as an
    editable Writer document instead, which is what a PDF->Word conversion
    actually needs.
    """
    # Give each invocation its own LO user profile so concurrent requests
    # (different jobs, different temp dirs) never collide on a lock file.
    profile_dir = os.path.join(out_dir, "_lo_profile")
    os.makedirs(profile_dir, exist_ok=True)
    profile_uri = f"file://{profile_dir}"

    out_subdir = os.path.join(out_dir, "_lo_output")
    os.makedirs(out_subdir, exist_ok=True)

    cmd = [
        SOFFICE_BIN,
        "--headless",
        "--norestore",
        "--nolockcheck",
        f"-env:UserInstallation={profile_uri}",
    ]
    if infilter:
        cmd += [f"--infilter={infilter}"]
    cmd += [
        "--convert-to",
        target_filter,
        "--outdir",
        out_subdir,
        input_path,
    ]
    try:
        result = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            timeout=CONVERT_TIMEOUT,
            text=True,
        )
    except subprocess.TimeoutExpired as exc:
        raise ConversionError(
            "The conversion took too long. The file may be very large or complex."
        ) from exc

    if result.returncode != 0:
        raise ConversionError(
            "LibreOffice could not convert this file. It may be corrupted, "
            "password-protected, or in an unsupported format."
        )

    # LibreOffice names the output after the input's basename with the new
    # extension; locate it rather than assuming the exact extension string.
    # Filter by the target extension too, defensively, in case the input
    # file's own copy ever ends up alongside it.
    stem = os.path.splitext(os.path.basename(input_path))[0]
    target_ext = target_filter.split(":", 1)[0].lstrip(".").lower()
    matches = [
        m for m in glob.glob(os.path.join(out_subdir, f"{stem}.*"))
        if os.path.isfile(m) and m.lower().endswith("." + target_ext)
    ]
    if not matches:
        raise ConversionError("Conversion finished but produced no output file.")
    return matches[0]


def _widen_columns_to_fit(xlsx_path: str) -> None:
    """Rewrite an .xlsx in place so every column is wide enough for its
    fullest cell value, and flag every sheet to export at fit-to-page-width.

    Why this exists: Excel/Calc only shows a cell's full text if the column
    is wide enough *or* the next cell is empty -- if the next cell has its
    own content, overflow text is silently clipped in the rendered output.
    A workbook's on-screen column widths are often left at their default
    (narrower than the data actually needs) because that clipping is easy
    to miss on screen; exporting "as-is" to a flat, uneditable PDF bakes it
    in permanently. This does the same thing a person would do by
    double-clicking a column border to auto-fit it, before the PDF export
    runs. Verified in testing 2026-08-27: without this, even an ordinary
    5-column staff-list spreadsheet lost characters from "Department" and
    names in the PDF; with it, nothing was clipped.
    """
    import openpyxl
    from openpyxl.utils import get_column_letter

    wb = openpyxl.load_workbook(xlsx_path)
    for ws in wb.worksheets:
        widths = {}
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is None:
                    continue
                length = len(str(cell.value))
                col = cell.column
                if length > widths.get(col, 0):
                    widths[col] = length
        for col, length in widths.items():
            ws.column_dimensions[get_column_letter(col)].width = min(max(length + 2, 8), 60)
        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 0
        ws.sheet_properties.pageSetUpPr.fitToPage = True
    wb.save(xlsx_path)


def convert_office_to_pdf(input_path: str, out_dir: str) -> str:
    """Word/Excel/PowerPoint (and ODT/RTF/TXT) -> PDF.

    Spreadsheets get extra treatment (see _widen_columns_to_fit and
    _SINGLE_PAGE_SHEETS_FILTER above) so wide or long sheets don't have
    columns or rows silently cut off -- the output PDF's page size grows to
    fit the data instead of the data being clipped to an arbitrary fixed
    page size. Word and PowerPoint files already carry their own fixed
    page/slide size, so they convert as-is. Font fidelity for all three
    depends on the server having a metric-compatible substitute installed
    for whatever font the document specifies (see the Dockerfile's font
    packages) -- LibreOffice cannot render a font it doesn't have.
    """
    ext = os.path.splitext(input_path)[1].lower().lstrip(".")
    if ext not in SPREADSHEET_EXTS:
        return _run_soffice(input_path, out_dir, "pdf")

    xlsx_path = input_path
    if ext == "xls":
        # openpyxl can only read/write .xlsx -- convert the legacy binary
        # format to .xlsx first (a lossless hop through LibreOffice) so the
        # column-width fix below can be applied to it too.
        xlsx_path = _run_soffice(input_path, out_dir, "xlsx:Calc MS Excel 2007 XML")

    try:
        _widen_columns_to_fit(xlsx_path)
    except Exception:
        # Best-effort: if the auto-fit pass fails for any reason (unusual
        # workbook structure, a formula openpyxl can't evaluate, etc.),
        # fall through and convert the sheet as-is rather than failing the
        # whole job over a cosmetic improvement.
        pass

    return _run_soffice(xlsx_path, out_dir, _SINGLE_PAGE_SHEETS_FILTER)


def convert_pdf_to_word(input_path: str, out_dir: str) -> str:
    """PDF -> editable .docx.

    Tries pdf2docx first. Verified in testing 2026-08-27 on a realistic
    document (headings, body paragraphs, a table) that LibreOffice's own
    PDF-import path (the fallback below) reconstructed the page as a set of
    duplicated, absolutely-positioned text frames with ZERO ordinary body
    paragraphs -- text that a search or copy-all can't find, and that can
    visually overlap depending on the viewer. pdf2docx instead rebuilds
    real flowing paragraphs and an actual editable table, recovering
    effectively all of the source text, and reads each run's font, size,
    weight and color directly off the PDF rather than guessing. Falls back
    to the LibreOffice path only if pdf2docx itself raises (e.g. on an
    unusual or corrupted PDF), so a pdf2docx failure doesn't hard-fail the
    tool. Known fidelity gap: a bulleted/numbered list in the source PDF
    comes back as a plain block of text rather than a native Word list --
    the words are intact, just not styled as a list.
    """
    out_path = os.path.join(out_dir, "converted_via_pdf2docx.docx")
    try:
        from pdf2docx import Converter

        cv = Converter(input_path)
        try:
            cv.convert(out_path)
        finally:
            cv.close()
        if os.path.exists(out_path) and os.path.getsize(out_path) > 0:
            return out_path
    except Exception:
        pass

    return _run_soffice(
        input_path, out_dir, "docx:MS Word 2007 XML", infilter="writer_pdf_import"
    )


def convert_pdf_to_ppt(input_path: str, out_dir: str) -> str:
    """PDF -> .pptx, one slide per page, each slide holding a full-page
    image of that page.

    Verified 2026-08-27: LibreOffice's headless --convert-to path (PDF
    opened as a Draw document, exported to "Impress MS PowerPoint 2007
    XML") silently produced an EMPTY .pptx -- zero slides, no slide master
    -- for every test PDF tried, including the simplest possible two-page
    text-only document, while still reporting success (exit code 0).
    Re-typing a Draw document as an Impress one doesn't work reliably via
    --convert-to, and going through Draw's own native format (.odg) as an
    intermediate step produced the same empty result -- the PDF imports
    into Draw correctly (confirmed via its content.xml), but nothing
    survives the Draw-to-Impress step.
    Given that, each page is rendered to an image (via pdf2image/poppler,
    the same renderer the PDF-to-JPG/PNG tools use) and placed as a single
    image sized to that page's exact original dimensions, with the overall
    slide size fixed to the first page's dimensions. This trades away
    editable text -- each slide is a picture, not text you can click into
    -- for a guaranteed, visually exact replica of every page (correct
    colors, fonts-as-drawn, spacing, everything), which is a more honest
    result than a broken "editable" file that silently has nothing in it.
    A true editable-text PDF->PPT converter is a separate, larger piece of
    work -- there is no equivalent to pdf2docx for PPTX that we're aware
    of; see the delivery notes for this trade-off.
    """
    import pypdf
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Emu

    try:
        reader = pypdf.PdfReader(input_path)
        if len(reader.pages) == 0:
            raise ConversionError("This PDF has no pages to convert.")

        first_box = reader.pages[0].mediabox
        slide_w = Emu(max(int(float(first_box.width) * 12700), 1))
        slide_h = Emu(max(int(float(first_box.height) * 12700), 1))

        prs = Presentation()
        prs.slide_width = slide_w
        prs.slide_height = slide_h
        blank_layout = prs.slide_layouts[6]

        images = convert_from_path(input_path, dpi=200)
        for i, img in enumerate(images):
            page_box = reader.pages[i].mediabox
            page_w = max(int(float(page_box.width) * 12700), 1)
            page_h = max(int(float(page_box.height) * 12700), 1)

            # Fit this page's image inside the fixed slide size, preserving
            # its own aspect ratio, in case pages aren't all the same size.
            scale = min(slide_w / page_w, slide_h / page_h)
            draw_w, draw_h = int(page_w * scale), int(page_h * scale)
            left, top = (slide_w - draw_w) // 2, (slide_h - draw_h) // 2

            slide = prs.slides.add_slide(blank_layout)
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            slide.shapes.add_picture(buf, left, top, width=draw_w, height=draw_h)

        out_path = os.path.join(out_dir, "converted.pptx")
        prs.save(out_path)
        return out_path
    except ConversionError:
        raise
    except Exception as exc:
        raise ConversionError(
            "Could not convert this PDF to a presentation. It may be "
            "corrupted or password-protected."
        ) from exc
